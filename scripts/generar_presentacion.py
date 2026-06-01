"""
Genera la presentación institucional de DenunciaBot en formato PowerPoint.

Salida: docs/DenunciaBot_Presentacion_Secretaria.pptx
Uso:    python scripts/generar_presentacion.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


# =========================================================================
# Paleta institucional
# =========================================================================

AZUL_INSTITUCIONAL = RGBColor(0x0B, 0x2E, 0x5C)
AZUL_CLARO = RGBColor(0x1E, 0x5A, 0x96)
GRIS_TEXTO = RGBColor(0x33, 0x33, 0x33)
GRIS_CLARO = RGBColor(0xF2, 0xF2, 0xF2)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
DORADO = RGBColor(0xC9, 0xA2, 0x27)
VERDE_OK = RGBColor(0x2E, 0x7D, 0x32)


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# =========================================================================
# Helpers
# =========================================================================

def _set_solid_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _add_rect(slide, x, y, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _set_solid_fill(shape, color)
    return shape


def _add_text(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = GRIS_TEXTO,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
    font: str = "Calibri",
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_bullets(
    slide,
    items: list[str],
    x,
    y,
    w,
    h,
    *,
    size: int = 16,
    color: RGBColor = GRIS_TEXTO,
    line_spacing: float = 1.2,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = color


def _add_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    x,
    y,
    w,
    h,
    *,
    header_color: RGBColor = AZUL_INSTITUCIONAL,
    row_zebra: RGBColor = GRIS_CLARO,
    font_size: int = 13,
):
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = table_shape.table

    # Header
    for j, htext in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = htext
        run.font.name = "Calibri"
        run.font.bold = True
        run.font.size = Pt(font_size + 1)
        run.font.color.rgb = BLANCO

    # Rows
    for i, row in enumerate(rows, start=1):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_zebra if i % 2 == 0 else BLANCO
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell_text
            run.font.name = "Calibri"
            run.font.size = Pt(font_size)
            run.font.color.rgb = GRIS_TEXTO
    return table


def _add_quote_box(slide, text: str, x, y, w, h, *, size: int = 16):
    """Caja con borde dorado a la izquierda — para citar mensajes del bot."""
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.1), h)
    _set_solid_fill(border, DORADO)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x + Inches(0.1), y, w - Inches(0.1), h
    )
    _set_solid_fill(bg, GRIS_CLARO)

    box = slide.shapes.add_textbox(
        x + Inches(0.3), y + Inches(0.15), w - Inches(0.5), h - Inches(0.3)
    )
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.2
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.color.rgb = GRIS_TEXTO
        run.font.italic = True


def _add_header_band(slide, title: str, subtitle: str | None = None):
    """Banda superior azul institucional con título."""
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.8), AZUL_INSTITUCIONAL)
    _add_text(
        slide,
        title,
        Inches(0.5),
        Inches(0.15),
        SLIDE_W - Inches(1),
        Inches(0.5),
        size=24,
        bold=True,
        color=BLANCO,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        _add_text(
            slide,
            subtitle,
            Inches(0.5),
            Inches(0.9),
            SLIDE_W - Inches(1),
            Inches(0.4),
            size=14,
            color=AZUL_CLARO,
        )


def _add_footer(slide, n: int, total: int):
    _add_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), AZUL_INSTITUCIONAL)
    _add_text(
        slide,
        "DenunciaBot — Secretaría General de Integridad Pública",
        Inches(0.4),
        SLIDE_H - Inches(0.32),
        Inches(8),
        Inches(0.3),
        size=10,
        color=BLANCO,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_text(
        slide,
        f"{n} / {total}",
        SLIDE_W - Inches(1.2),
        SLIDE_H - Inches(0.32),
        Inches(0.8),
        Inches(0.3),
        size=10,
        color=BLANCO,
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.RIGHT,
    )


# =========================================================================
# Slides
# =========================================================================

def slide_portada(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, AZUL_INSTITUCIONAL)
    _add_rect(slide, 0, Inches(3.0), SLIDE_W, Inches(0.05), DORADO)

    _add_text(
        slide,
        "DenunciaBot",
        Inches(0.5),
        Inches(2.2),
        SLIDE_W - Inches(1),
        Inches(1),
        size=60,
        bold=True,
        color=BLANCO,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Chatbot institucional de denuncias ciudadanas de corrupción",
        Inches(0.5),
        Inches(3.3),
        SLIDE_W - Inches(1),
        Inches(0.6),
        size=24,
        color=BLANCO,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Secretaría General de Integridad Pública del Ecuador\nConvenio interinstitucional con SERCOP",
        Inches(0.5),
        Inches(4.5),
        SLIDE_W - Inches(1),
        Inches(1),
        size=18,
        color=DORADO,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Quito · Mayo de 2026",
        Inches(0.5),
        Inches(6.5),
        SLIDE_W - Inches(1),
        Inches(0.4),
        size=14,
        color=BLANCO,
        align=PP_ALIGN.CENTER,
    )


def slide_que_es(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, "¿Qué es DenunciaBot?")
    _add_text(
        slide,
        "Un canal conversacional por WhatsApp para que el ciudadano presente "
        "denuncias de corrupción de manera estructurada, confidencial y trazable.",
        Inches(0.5),
        Inches(1.4),
        SLIDE_W - Inches(1),
        Inches(1.2),
        size=18,
        color=GRIS_TEXTO,
    )
    _add_bullets(
        slide,
        [
            "Confidencial — el número telefónico nunca se almacena en claro.",
            "Estructurado — guiado por una secuencia fija de 8 preguntas.",
            "Determinista — sin inteligencia artificial generativa; respuestas controladas y trazables.",
            "Trazable — cada denuncia recibe un código único de seguimiento.",
        ],
        Inches(0.7),
        Inches(2.8),
        SLIDE_W - Inches(1.5),
        Inches(2.5),
        size=18,
    )
    _add_quote_box(
        slide,
        "El ciudadano abre WhatsApp, escribe al número institucional, "
        "y en menos de 5 minutos su denuncia queda registrada y notificada "
        "al equipo de la Secretaría.",
        Inches(0.5),
        Inches(5.5),
        SLIDE_W - Inches(1),
        Inches(1.2),
        size=16,
    )
    _add_footer(slide, n, total)


def slide_principios(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, "Principios no negociables")
    _add_table(
        slide,
        ["#", "Principio", "Garantía"],
        [
            ["1", "Sin IA generativa", "Flujo 100% determinista, auditable, sin alucinaciones"],
            ["2", "Cifrado de campos sensibles", "Institución, descripción y nombres se cifran antes de guardarse"],
            ["3", "Teléfono nunca en claro", "Solo se almacena un hash irreversible con sal secreta"],
            ["4", "Bitácora inmutable", "Ningún cambio en el registro puede borrarse — garantía legal"],
            ["5", "Validación de origen", "Cada mensaje entrante se verifica criptográficamente"],
            ["6", "Escaneo antivirus", "Las evidencias adjuntas se inspeccionan antes de almacenarse"],
        ],
        Inches(0.5),
        Inches(1.5),
        SLIDE_W - Inches(1),
        Inches(5),
        font_size=14,
    )
    _add_footer(slide, n, total)


def slide_recorrido(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, "Recorrido del ciudadano")
    pasos = [
        "1.  Saludo y aceptación de condiciones",
        "2.  Pregunta 1 — Institución denunciada",
        "3.  Pregunta 2 — Descripción de los hechos",
        "4.  Pregunta 3 — Fecha aproximada",
        "5.  Pregunta 4 — Personas involucradas (opcional)",
        "6.  Pregunta 5 — Perjuicio económico (opcional)",
        "7.  Pregunta 6 — Denuncia previa en otra entidad (opcional)",
        "8.  Pregunta 7 — Evidencias (PDF/JPG/PNG, opcional)",
        "9.  Pregunta 8 — Resumen y confirmación",
        "10. Código de seguimiento entregado al ciudadano",
        "11. Notificación por correo al equipo de la Secretaría",
    ]
    box = slide.shapes.add_textbox(
        Inches(1), Inches(1.4), SLIDE_W - Inches(2), Inches(5.5)
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, paso in enumerate(pasos):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = paso
        run.font.name = "Consolas"
        run.font.size = Pt(16)
        run.font.color.rgb = GRIS_TEXTO
    _add_footer(slide, n, total)


def _slide_pregunta(prs, n, total, titulo: str, mensaje: str, validacion: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, titulo, "El bot envía al ciudadano:")
    _add_quote_box(
        slide,
        mensaje,
        Inches(0.5),
        Inches(1.6),
        SLIDE_W - Inches(1),
        Inches(4.3),
        size=15,
    )
    _add_text(
        slide,
        "Validación:",
        Inches(0.5),
        Inches(6.1),
        Inches(2),
        Inches(0.4),
        size=14,
        bold=True,
        color=AZUL_INSTITUCIONAL,
    )
    _add_text(
        slide,
        validacion,
        Inches(1.8),
        Inches(6.1),
        SLIDE_W - Inches(2.3),
        Inches(0.8),
        size=14,
        color=GRIS_TEXTO,
    )
    _add_footer(slide, n, total)


def slide_bienvenida(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, "Inicio: bienvenida y aceptación de condiciones")
    _add_quote_box(
        slide,
        "Bienvenido al sistema institucional de denuncias de corrupción.\n\n"
        "Este canal pertenece a la Secretaría General de Integridad Pública "
        "del Ecuador. A través de este chat usted podrá registrar de forma "
        "confidencial una denuncia ciudadana sobre presuntos hechos de "
        "corrupción en el sector público.\n\n"
        "•  Su número de teléfono no se almacena en claro.\n"
        "•  Su denuncia es confidencial y se cifra antes de almacenarse.\n"
        "•  Al finalizar recibirá un código de seguimiento único.\n"
        "•  Puede escribir cancelar en cualquier momento.\n"
        "•  Si no responde durante 5 minutos la sesión se cierra.",
        Inches(0.5),
        Inches(1.4),
        SLIDE_W - Inches(1),
        Inches(4.5),
        size=13,
    )
    _add_text(
        slide,
        "Consentimiento explícito:",
        Inches(0.5),
        Inches(6.1),
        Inches(3.5),
        Inches(0.4),
        size=14,
        bold=True,
        color=AZUL_INSTITUCIONAL,
    )
    _add_text(
        slide,
        "¿Acepta continuar?   [ Sí ]   [ No ]",
        Inches(3.5),
        Inches(6.1),
        SLIDE_W - Inches(4),
        Inches(0.4),
        size=14,
        color=GRIS_TEXTO,
    )
    _add_footer(slide, n, total)


def slide_resumen(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, "Pregunta 8: resumen y confirmación")
    _add_quote_box(
        slide,
        "Paso 8 de 8 — Resumen de su denuncia\n\n"
        "Institución: Ministerio de Salud Pública\n"
        "Descripción: [texto resumido…]\n"
        "Fecha aproximada: 15/03/2025\n"
        "Personas involucradas: Juan Pérez, director administrativo\n"
        "Perjuicio económico: aprox. 50 000 USD\n"
        "Denuncia previa: No\n"
        "Evidencias adjuntas: 2\n\n"
        "Si confirma, registraremos su denuncia y le entregaremos un código "
        "de seguimiento. Si desea modificar algún dato, pulse Editar.",
        Inches(0.5),
        Inches(1.4),
        SLIDE_W - Inches(1),
        Inches(4.8),
        size=13,
    )
    _add_text(
        slide,
        "Opciones del ciudadano:   [ Confirmar ]   [ Editar ]   [ Cancelar ]",
        Inches(0.5),
        Inches(6.4),
        SLIDE_W - Inches(1),
        Inches(0.4),
        size=14,
        bold=True,
        color=AZUL_INSTITUCIONAL,
        align=PP_ALIGN.CENTER,
    )
    _add_footer(slide, n, total)


def slide_cierre_codigo(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, "Cierre: entrega del código de seguimiento")
    _add_quote_box(
        slide,
        "Su denuncia ha sido registrada exitosamente.\n\n"
        "Código de seguimiento:  ALR-2026-MJY3LW\n\n"
        "Conserve este código: le permitirá consultar el estado de su "
        "denuncia con la entidad. La información ha sido cifrada y enviada "
        "a la unidad correspondiente para su revisión.\n\n"
        "Gracias por colaborar con la integridad pública del Ecuador.",
        Inches(0.5),
        Inches(1.4),
        SLIDE_W - Inches(1),
        Inches(3.8),
        size=16,
    )
    _add_text(
        slide,
        "Consulta pública del estado de la denuncia:",
        Inches(0.5),
        Inches(5.6),
        SLIDE_W - Inches(1),
        Inches(0.4),
        size=14,
        bold=True,
        color=AZUL_INSTITUCIONAL,
    )
    _add_text(
        slide,
        "https://denuncia.secretaria.gob.ec/alerta/ALR-2026-MJY3LW",
        Inches(0.5),
        Inches(6.0),
        SLIDE_W - Inches(1),
        Inches(0.4),
        size=15,
        color=AZUL_CLARO,
        font="Consolas",
    )
    _add_footer(slide, n, total)


def slide_que_pasa_despues(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, "¿Qué pasa después de la confirmación?")
    _add_bullets(
        slide,
        [
            "La denuncia queda cifrada en la base de datos institucional.",
            "Se genera un código público único (alfabeto sin caracteres ambiguos).",
            "Se registra el evento en la bitácora inmutable.",
            "Se envía un correo automático al buzón integridad@secretaria.gob.ec con el detalle para revisión humana.",
            "El ciudadano recibe su código por WhatsApp antes de 5 segundos desde la confirmación.",
        ],
        Inches(0.6),
        Inches(1.6),
        SLIDE_W - Inches(1.2),
        Inches(5),
        size=18,
    )
    _add_footer(slide, n, total)


def slide_panel(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, "Panel administrativo")
    _add_text(
        slide,
        "Panel web minimalista para el equipo de la Secretaría:",
        Inches(0.5),
        Inches(1.4),
        SLIDE_W - Inches(1),
        Inches(0.5),
        size=18,
        color=GRIS_TEXTO,
    )
    _add_bullets(
        slide,
        [
            "Listado de todas las denuncias recibidas, con paginación.",
            "Filtros por estado: REGISTRADA, EN_REVISION, TRAMITADA, DESCARTADA.",
            "Detalle de cada denuncia con datos descifrados al vuelo.",
            "Cambio de estado auditado — cada cambio queda en la bitácora con la identidad del operador.",
            "Exportación de la bitácora completa con firma criptográfica, para fines legales.",
        ],
        Inches(0.7),
        Inches(2.2),
        SLIDE_W - Inches(1.5),
        Inches(3.5),
        size=16,
    )
    _add_quote_box(
        slide,
        "Acceso protegido por token compartido. "
        "La sesión expira automáticamente tras 8 horas.",
        Inches(0.5),
        Inches(5.8),
        SLIDE_W - Inches(1),
        Inches(0.8),
        size=15,
    )
    _add_footer(slide, n, total)


def slide_proteccion(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, "Mecanismos de protección al ciudadano")
    _add_table(
        slide,
        ["Situación", "Respuesta del bot"],
        [
            ["Ciudadano escribe 'cancelar'", "Aborta el flujo sin guardar nada"],
            ["4 minutos sin responder", "Aviso de inactividad"],
            ["5 minutos sin responder", "Cierre automático; sesión descartada"],
            ["Respuesta inválida", "3 reintentos con orientación, luego cierre limpio"],
            ["Archivo malicioso", "Antivirus rechaza; ciudadano sigue el flujo"],
            ["Saludo repetido tras pausa", "Retoma la conversación donde quedó"],
        ],
        Inches(0.5),
        Inches(1.5),
        SLIDE_W - Inches(1),
        Inches(5),
        font_size=15,
    )
    _add_footer(slide, n, total)


def slide_estado_actual(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, "Estado actual del proyecto", "Construido y validado a mayo de 2026")
    _add_bullets(
        slide,
        [
            "Flujo completo de las 8 preguntas (12 estados internos).",
            "Cifrado de campos sensibles y hash de teléfono con sal secreta.",
            "Bitácora inmutable a nivel de base de datos.",
            "Webhook con validación criptográfica de cada mensaje entrante de Meta.",
            "Panel administrativo web funcional.",
            "Endpoint público de consulta por código de seguimiento.",
            "Notificación automática por correo institucional.",
            "Respaldos automatizados de la base de datos.",
            "Pruebas automatizadas: 171 pruebas unitarias en verde.",
            "Manual de operación documentado.",
        ],
        Inches(0.7),
        Inches(1.6),
        SLIDE_W - Inches(1.5),
        Inches(5),
        size=16,
    )
    _add_footer(slide, n, total)


def slide_faltantes(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, "Lo que falta para producción")
    _add_table(
        slide,
        ["Requisito", "Responsable", "Estado"],
        [
            ["Casilla institucional denunciabot@sercop.gob.ec", "IT SERCOP", "Solicitada"],
            ["Servidor RHEL con Python 3.11 y Redis", "IT SERCOP", "Solicitada"],
            ["Subdominio público + certificado TLS", "Seguridad SERCOP", "Solicitada"],
            ["Número de WhatsApp Business certificado", "Meta + Secretaría", "Por gestionar"],
            ["Capacitación al equipo de revisión", "Secretaría", "Por agendar"],
            ["Procedimiento interno de revisión", "Secretaría", "Por definir"],
        ],
        Inches(0.5),
        Inches(1.5),
        SLIDE_W - Inches(1),
        Inches(5),
        font_size=14,
    )
    _add_footer(slide, n, total)


def slide_cronograma(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header_band(slide, "Cronograma estimado", "Ajustable según tiempos institucionales")
    _add_table(
        slide,
        ["Semana", "Hito"],
        [
            ["1", "IT entrega infraestructura base"],
            ["2", "Despliegue en servidor institucional + pruebas internas"],
            ["3", "Certificación del número de WhatsApp con Meta"],
            ["4", "Capacitación al equipo de la Secretaría"],
            ["5", "Marcha blanca con casos controlados"],
            ["6", "Lanzamiento público con la Secretaría"],
        ],
        Inches(0.5),
        Inches(1.5),
        SLIDE_W - Inches(1),
        Inches(5),
        font_size=15,
    )
    _add_footer(slide, n, total)


def slide_cierre(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, AZUL_INSTITUCIONAL)
    _add_rect(slide, 0, Inches(3.5), SLIDE_W, Inches(0.05), DORADO)
    _add_text(
        slide,
        "DenunciaBot está técnicamente listo.",
        Inches(0.5),
        Inches(2.4),
        SLIDE_W - Inches(1),
        Inches(0.8),
        size=36,
        bold=True,
        color=BLANCO,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Lo que sigue es gestión institucional:\ninfraestructura, número de WhatsApp y procedimiento interno de revisión.",
        Inches(0.5),
        Inches(3.7),
        SLIDE_W - Inches(1),
        Inches(1.2),
        size=20,
        color=DORADO,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "¿Preguntas?",
        Inches(0.5),
        Inches(5.0),
        SLIDE_W - Inches(1),
        Inches(0.6),
        size=28,
        bold=True,
        color=BLANCO,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        "Jonathan Mauricio Ruiz Sánchez\nAnalista de Operaciones de Innovación Tecnológica 2 — SERCOP",
        Inches(0.5),
        Inches(6.3),
        SLIDE_W - Inches(1),
        Inches(0.8),
        size=14,
        color=BLANCO,
        align=PP_ALIGN.CENTER,
    )


# =========================================================================
# Main
# =========================================================================

PREGUNTAS = [
    (
        "Pregunta 1: Institución denunciada",
        "Paso 1 de 8. Indique el nombre de la institución, dependencia o "
        "entidad pública sobre la cual presenta la denuncia.\n\n"
        "Ejemplos: \"Ministerio de Salud Pública\", \"GAD Municipal de "
        "Quito\", \"Hospital del Seguro Social en Guayaquil\".",
        "Entre 3 y 200 caracteres. Se cifra antes de almacenarse.",
    ),
    (
        "Pregunta 2: Descripción de los hechos",
        "Paso 2 de 8. Describa de manera detallada los hechos que desea "
        "denunciar. Incluya, si los conoce: lugar, modalidad del presunto "
        "acto, montos involucrados, fechas relevantes y todo dato que "
        "considere útil para la investigación.\n\n"
        "Mínimo 30 caracteres, máximo 2000.",
        "Entre 30 y 2000 caracteres. Se cifra antes de almacenarse.",
    ),
    (
        "Pregunta 3: Fecha aproximada",
        "Paso 3 de 8. Indique la fecha aproximada en que ocurrieron los hechos.\n\n"
        "Formatos válidos:\n"
        "•  Día/mes/año (ej. 15/03/2025)\n"
        "•  Mes/año (ej. 03/2025)\n"
        "•  Solo año (ej. 2025)\n"
        "•  La frase \"no recuerdo\"",
        "Fecha plausible (no futura, no anterior a 1990) o \"no recuerdo\".",
    ),
    (
        "Pregunta 4: Personas involucradas  (opcional)",
        "Paso 4 de 8. Si conoce a las personas presuntamente involucradas, "
        "indique nombres, cargos o cualquier referencia que pueda ayudar.\n\n"
        "Si no las conoce, escriba \"no conozco\" o \"prefiero no decir\".",
        "Texto libre o frase de exclusión. Si hay datos, se cifran.",
    ),
    (
        "Pregunta 5: Perjuicio económico  (opcional)",
        "Paso 5 de 8. Si conoce el monto aproximado del perjuicio "
        "económico, indíquelo (en dólares).\n\n"
        "Puede escribir un número, un rango (\"entre 5000 y 10000\") o "
        "\"no aplica\" / \"no sé\".",
        "Número, rango o frase de exclusión.",
    ),
    (
        "Pregunta 6: Denuncia previa en otra entidad  (opcional)",
        "Paso 6 de 8. ¿Ha presentado anteriormente esta misma denuncia "
        "ante otra entidad (Fiscalía, Contraloría, otra)?\n\n"
        "[ Sí ]   [ No ]\n\n"
        "Si responde Sí, el bot pregunta:\n"
        "\"Indique brevemente ante qué entidad presentó la denuncia previa "
        "y, si lo recuerda, cuándo lo hizo.\"",
        "Botones Sí/No. Si Sí: texto libre.",
    ),
    (
        "Pregunta 7: Evidencias  (opcional)",
        "Paso 7 de 8. Si dispone de evidencias, envíelas ahora.\n\n"
        "•  Formatos aceptados: PDF, JPG, PNG\n"
        "•  Tamaño máximo por archivo: 10 MB\n"
        "•  Cantidad máxima: 5 archivos\n\n"
        "Envíe los archivos uno por uno. Cuando termine, pulse \"No tengo "
        "más\" o escriba \"no tengo\".",
        "Tipo MIME validado + escaneo antivirus. Nombre original cifrado.",
    ),
]


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 21

    slide_portada(prs)
    slide_que_es(prs, 2, total)
    slide_principios(prs, 3, total)
    slide_recorrido(prs, 4, total)
    slide_bienvenida(prs, 5, total)

    for i, (titulo, mensaje, validacion) in enumerate(PREGUNTAS, start=6):
        _slide_pregunta(prs, i, total, titulo, mensaje, validacion)

    slide_resumen(prs, 13, total)
    slide_cierre_codigo(prs, 14, total)
    slide_que_pasa_despues(prs, 15, total)
    slide_panel(prs, 16, total)
    slide_proteccion(prs, 17, total)
    slide_estado_actual(prs, 18, total)
    slide_faltantes(prs, 19, total)
    slide_cronograma(prs, 20, total)
    slide_cierre(prs, 21, total)

    out_dir = Path(__file__).resolve().parent.parent / "docs"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "DenunciaBot_Presentacion_Secretaria.pptx"
    prs.save(str(out_path))
    print(f"OK  Presentación generada: {out_path}")
    print(f"    Total de diapositivas: {len(prs.slides)}")


if __name__ == "__main__":
    main()
